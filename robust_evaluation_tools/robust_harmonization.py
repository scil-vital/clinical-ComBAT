import pandas as pd
import numpy as np
import subprocess
import os
import re

from scripts import combat_quick_apply
from scripts import combat_quick_QC
from robust_evaluation_tools.robust_utils import get_site, robust_text, rwp_text

from pptx import Presentation
from pptx.util import Inches

def get_output_model_filename(mov_data_file, metric, harmonizartion_method, robust, rwp):
        return (
            get_site(mov_data_file)
            + "."
            + metric
            + "."
            + harmonizartion_method
            + "."
            + robust_text(robust)
            + "."
            + rwp_text(rwp)
            + ".model.csv"
        )

def get_output_filename(mov_data_file, metric, harmonizartion_method, robust, rwp, directory):
        site = get_site(mov_data_file)
        if "test" in mov_data_file:
            site += "_test"
        return os.path.join(
            directory,
            site
            + "."
            + metric
            + "."
            + harmonizartion_method
            + "."
            + robust_text(robust)
            + "."
            + rwp_text(rwp)
            + ".csv"
        )

def fit(mov_data_file, ref_data_file, metric, harmonizartion_method, robust, rwp, directory, hc,):
    ###########
    ### fit ###
    ###########
    
    output_model_filename = get_output_model_filename(
        mov_data_file, metric, harmonizartion_method, robust, rwp
    )
        # Check if the output model file already exists
    output_model_path = os.path.join(directory, output_model_filename)
    if os.path.exists(output_model_path):
        return output_model_path
    cmd = (
        "scripts/combat_quick_fit.py"
        + " "
        + ref_data_file
        + " "
        + mov_data_file
        + " --out_dir "
        + directory
        + " --output_model_filename "
        + output_model_filename
        + " --method "
        + harmonizartion_method
        + " --robust "
        + robust
        + " -f "
        + " --no_empirical_bayes"
    )
    if rwp:
        cmd += ' --rwp'
    if hc: 
        cmd += ' --hc'

    subprocess.call(cmd, shell=True)
    return output_model_path

def apply(mov_data_file, model_filename, metric, harmonizartion_method, robust, rwp, directory):
    output_filename = get_output_filename(mov_data_file, metric, harmonizartion_method, robust, rwp, directory)
    if os.path.exists(output_filename):
        return output_filename
    combat_quick_apply.apply(mov_data_file, model_filename, output_filename)
    return output_filename

def visualize_harmonization(f, new_f, ref_data_file, directory, bundles = '', title=''):
    cmd = (
        "scripts/combat_visualize_harmonization.py"
        + " "
        + ref_data_file
        + " "
        + f
        + " "
        + new_f
        + " --out_dir "
        + directory
        + " -f"
    )
    if bundles != '':
        cmd += f" --bundles {bundles}"
    if title != '':
        cmd += f" --outname {title}"
        cmd += f" --add_suffix {title}"
    subprocess.call(cmd, shell=True)

def QC(ref_data, output_filename, output_model_filename):
    return combat_quick_QC.QC(ref_data, output_filename, output_model_filename)


def compare_with_compilation_STD(df, compilation_df):

    # Filtrer les patients de COMPILATION qui sont dans df en utilisant les sid
    common_sids = df['sid'].unique()
    filtered_compilation_df = compilation_df[compilation_df['sid'].isin(common_sids)]

    if len(filtered_compilation_df) != len(df):
        raise ValueError(f"Attention: Nombre de lignes différent entre df ({len(df)}) et filtered_compilation_df ({len(filtered_compilation_df)})")

    # Initialiser une liste pour stocker les résultats
    comparison_df = pd.DataFrame()

    # Comparer la différence absolue de la colonne mean par bundle
    for bundle in df['bundle'].unique():
        df_bundle = df[df['bundle'] == bundle]
        compilation_bundle = filtered_compilation_df[filtered_compilation_df['bundle'] == bundle]
        std_val = compilation_bundle['mean_no_cov'].std()
        
        # Fusionner les deux DataFrames sur les colonnes 'sid' et 'bundle'
        merged_df = pd.merge(df_bundle, compilation_bundle, on=['sid', 'bundle'], suffixes=('_df', '_compilation'))
        
        # Calculer la différence absolue de la colonne mean
        merged_df['abs_diff_mean'] = (merged_df['mean_df'] - merged_df['mean_compilation']).abs() / std_val
        # Calculer la somme des différences absolues pour le bundle
        comparison_df[bundle] = merged_df['abs_diff_mean']
            
    # Ajouter le site au DataFrame
    mean_df = pd.DataFrame(comparison_df.mean()).transpose()
    return mean_df
    
def compare_with_compilation(df, compilation_df):

    # Filtrer les patients de COMPILATION qui sont dans df en utilisant les sid
    common_sids = df['sid'].unique()
    filtered_compilation_df = compilation_df[compilation_df['sid'].isin(common_sids)]

    if len(filtered_compilation_df) != len(df):
        raise ValueError(f"Attention: Nombre de lignes différent entre df ({len(df)}) et filtered_compilation_df ({len(filtered_compilation_df)})")

    # Initialiser une liste pour stocker les résultats
    comparison_df = pd.DataFrame()

    # Comparer la différence absolue de la colonne mean par bundle
    for bundle in df['bundle'].unique():
        df_bundle = df[df['bundle'] == bundle]
        compilation_bundle = filtered_compilation_df[filtered_compilation_df['bundle'] == bundle]
        
        # Fusionner les deux DataFrames sur les colonnes 'sid' et 'bundle'
        merged_df = pd.merge(df_bundle, compilation_bundle, on=['sid', 'bundle'], suffixes=('_df', '_compilation'))
        
        # Calculer la différence absolue de la colonne mean
        merged_df['abs_diff_mean'] = (merged_df['mean_df'] - merged_df['mean_compilation']).abs()
        # Calculer la somme des différences absolues pour le bundle
        comparison_df[bundle] = merged_df['abs_diff_mean']
            
    # Ajouter le site au DataFrame
    mean_df = pd.DataFrame(comparison_df.mean()).transpose()

    return mean_df

def compare_with_compilation_SMAPE(df, compilation_df):
    # compilation_df = get_ground_truth_df(df, directory, harmonizartion_method)
    common_sids = df['sid'].unique()
    comp_filt = compilation_df[compilation_df['sid'].isin(common_sids)]

    if len(comp_filt) != len(df):
        raise ValueError(f"Attention: Nombre de lignes différent entre df ({len(df)}) et filtered_compilation_df ({len(comp_filt)})")

    def smape(y_true, y_pred):
        denom = (np.abs(y_true) + np.abs(y_pred)) / 2
        return np.where(denom == 0, 0, np.abs(y_true - y_pred) / denom) * 100

    comparison_df = pd.DataFrame()
    for bundle in df['bundle'].unique():
        merged = pd.merge(
            df[df['bundle'] == bundle],
            comp_filt[comp_filt['bundle'] == bundle],
            on=['sid', 'bundle'],
            suffixes=('_df', '_compilation')
        )
        comparison_df[bundle] = smape(
            merged['mean_df'].to_numpy(),
            merged['mean_compilation'].to_numpy()
        )

    mean_df = pd.DataFrame(comparison_df.mean()).T
    mean_df.index = ['SMAPE_mean_%']
    return mean_df

def compare_with_compilation_var(df, compilation_df):

    # Filtrer les patients de COMPILATION qui sont dans df en utilisant les sid
    common_sids = df['sid'].unique()
    filtered_compilation_df = compilation_df[compilation_df['sid'].isin(common_sids)]

    if len(filtered_compilation_df) != len(df):
        raise ValueError(f"Attention: Nombre de lignes différent entre df ({len(df)}) et filtered_compilation_df ({len(filtered_compilation_df)})")

    # Initialiser une liste pour stocker les résultats
    comparison_df = pd.DataFrame()

    # Comparer la différence de variance par bundle
    for bundle in df['bundle'].unique():
        df_bundle = df[df['bundle'] == bundle]
        compilation_bundle = filtered_compilation_df[filtered_compilation_df['bundle'] == bundle]
        
        # Calculer la variance pour chaque bundle
        variance_df = df_bundle['mean'].var()
        variance_compilation = compilation_bundle['mean'].var()
        
        # Calculer la différence absolue des variances
        comparison_df[bundle] = [abs(variance_df - variance_compilation)]
            
    # Ajouter le site au DataFrame
    mean_df = pd.DataFrame(comparison_df.mean()).transpose()

    return mean_df

def compare_with_camcan(df):
    compilation_df = get_camcan_df(df)
    # Charger le DataFrame COMPILATION

    # Filtrer les patients de COMPILATION qui sont dans df en utilisant les sid
    common_sids = df['sid'].unique()
    filtered_compilation_df = compilation_df[compilation_df['sid'].isin(common_sids)]

    # Initialiser une liste pour stocker les résultats
    comparison_df = pd.DataFrame()

    # Comparer la différence absolue de la colonne mean par bundle
    for bundle in df['bundle'].unique():
        df_bundle = df[df['bundle'] == bundle]
        compilation_bundle = filtered_compilation_df[filtered_compilation_df['bundle'] == bundle]
        
        # Fusionner les deux DataFrames sur les colonnes 'sid' et 'bundle'
        merged_df = pd.merge(df_bundle, compilation_bundle, on=['sid', 'bundle'], suffixes=('_df', '_compilation'))
        
        # Calculer la différence absolue de la colonne mean
        merged_df['abs_diff_mean'] = (merged_df['mean_df'] - merged_df['mean_compilation']).abs()
        # Calculer la somme des différences absolues pour le bundle
        comparison_df[bundle] = merged_df['abs_diff_mean']
            
    # Ajouter le site au DataFrame
    mean_df = pd.DataFrame(comparison_df.mean()).transpose()

    return mean_df


def create_presentation(directory, method):
    # Create a presentation object
    prs = Presentation()
    
    # Define the subdirectories
    subdirs = ["hc", "NoRobust", "robust", "robust_rwp"]
    # Get the list of images
    images = [img for img in os.listdir(os.path.join(directory, subdirs[0])) if method in img and img.endswith('.png')]
    
    for img in images:
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        for i, subdir in enumerate(subdirs):
            img_path = os.path.join(directory, subdir, img)
            left = Inches(0.5 + (i % 2) * 4.5)  # Positioning images in two columns
            top = Inches(0.2 + (i // 2) * 3.5)  # Positioning images in two rows with more space between rows
            
            # Add text above the image
            text_box = slide.shapes.add_textbox(left, top, width=Inches(4), height=Inches(0.5))
            text_frame = text_box.text_frame
            text_frame.text = subdir
            
            # Add the image
            slide.shapes.add_picture(img_path, left, top + Inches(0.5), width=Inches(4))
    
    # Save the presentation
    prs.save(os.path.join(directory, 'harmonization_results.pptx'))


def compare_distances(directory, site, hc_dists, no_robust_dists, robust_dists, robust_rwp_dists):
    # compare les distances de 4 methodes de harmonization
    comparison_results = {
        "hc_vs_no_robust": (np.array(hc_dists) - np.array(no_robust_dists))/np.array(no_robust_dists)*100,
        "robust_vs_no_robust": (np.array(robust_dists) - np.array(no_robust_dists))/np.array(no_robust_dists)*100,
        "robust_rwp_vs_no_robust": (np.array(robust_rwp_dists) - np.array(no_robust_dists))/np.array(no_robust_dists)*100
    }
    df = pd.DataFrame(comparison_results)
    
    # Calculer le nombre de comparaisons négatives et positives, et les moyennes et médianes
    results = []
    for method in comparison_results.keys():
        negative_values = df[method][df[method] < 0]
        positive_values = df[method][df[method] >= 0]
        
        num_negative = len(negative_values)
        num_positive = len(positive_values)
        
        mean_negative = negative_values.mean() if num_negative > 0 else 0
        mean_positive = positive_values.mean() if num_positive > 0 else 0
        
        median_negative = negative_values.median() if num_negative > 0 else 0
        median_positive = positive_values.median() if num_positive > 0 else 0
        
        mean_difference = df[method].mean()
        
        results.append({
            "site": site,
            "comparaison": method,
            "Nb comp. nég.": num_negative,
            "Nb comp. pos.": num_positive,
            "Moy. tot.": mean_difference,
            "Moy. val. nég.": mean_negative,
            "Moy. val. pos.": mean_positive,
            "Méd. val. nég.": median_negative,
            "Méd. val. pos.": median_positive
        })
    results_df = pd.DataFrame(results)
    results_df.to_csv(os.path.join(directory, f"{site}_comparison_results.csv"), index=False)
    return results_df


def get_csv_res(mov_data_file, directory, harmonizartion_method,metric):
    mov_data = pd.read_csv(mov_data_file)
    model = (
            directory
            +
            '/'
            +
            mov_data.site.unique()[0]
            + "."
            + metric
            + "."
            + harmonizartion_method
            + ".model.csv"
        )
    return combat_quick_apply.make_best(mov_data_file, model)


def get_compilation_df(df):
    disease = df['site'].iloc[0].split('_', 1)[0]
    metric = df['metric'].unique()[0]

    compilation_folder = os.path.join('DONNES_F', 'COMPILATIONS_AUG_3')
    compilation_file = os.path.join(compilation_folder, f"{disease}_combination_all_metrics_CamCAN.csv.gz")
    compilation_df = pd.read_csv(compilation_file, compression='gzip')
    compilation_df = compilation_df[~compilation_df['bundle'].isin(['left_ventricle', 'right_ventricle'])]
    return compilation_df[compilation_df['metric'] == metric]

def get_ground_truth_df(df, directory, harmonizartion_method):
    metric = df['metric'].unique()[0]
    gt_folder = os.path.join(directory, 'hc')
    compilation_file = os.path.join(
            gt_folder,
            str(df.site.unique()[0])
            + "."
            + metric
            + "."
            + harmonizartion_method
            + "."
            + robust_text("No")
            + "."
            + rwp_text(False)
            + ".csv"
        )
    compilation_df = pd.read_csv(compilation_file)
    return compilation_df

def get_camcan_df(df):
    metric = df['metric'].unique()[0]

    compilation_folder = os.path.join('DONNES_F', 'CamCAN')
    camcan_file = os.path.join(compilation_folder, f"CamCAN.{metric}.raw.csv.gz")
    compilation_df = pd.read_csv(camcan_file, compression='gzip')
    return compilation_df
